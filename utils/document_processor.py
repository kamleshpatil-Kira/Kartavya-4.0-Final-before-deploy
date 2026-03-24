from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from pathlib import Path
import zipfile
import shutil
from typing import Dict, Any, List, Optional


def _is_text_garbled(text: str, min_words: int = 50, max_null_ratio: float = 0.15) -> bool:
    """Return True if extracted text is too short, contains too many null/replacement chars,
    or has an unnaturally high frequency of a single character/block.
    Detects Type3 / custom-encoded PDFs where glyph->unicode mapping is absent (common in Hindi PDFs)."""
    if not text or not text.strip():
        return True
    
    non_space_chars = len([c for c in text if not c.isspace()])
    if non_space_chars == 0:
        return True
        
    # Standard valid unicode letter check (works for Devanagari, Latin, Arabic, etc.)
    alpha_chars = len([c for c in text if c.isalpha()])
    
    # If less than 20% of the document is made up of actual letters (e.g. 90% is '■' or math symbols)
    if alpha_chars / non_space_chars < 0.20:
        return True

    words = text.split()
    if len(words) < min_words:
        return True
        
    null_chars = text.count(chr(0)) + text.count(chr(0xfffd)) + text.count('■') + text.count('')
    null_ratio = null_chars / max(len(text), 1)
    if null_ratio > max_null_ratio:
        return True
        
    # Detect repeating-character garble (e.g. all unmapped glyphs extract as 'n's)
    import collections
    letters = [c.lower() for c in text if c.isalpha()]
    if letters:
        most_common_count = collections.Counter(letters).most_common(1)[0][1]
        if (most_common_count / len(letters)) > 0.40:  # If one letter is >40% of all letters, it's garbled
            return True
            
    return False

from utils.logger import logger, log_activity
import tempfile
import re
import os

class DocumentProcessor:
    """Process uploaded documents (.docx, .pdf, .pptx, .zip, audio, video)"""
    
    def __init__(self):
        # Initialize Whisper model (load once, reuse)
        self.whisper_model = None
    
    def _get_whisper_model(self):
        """Lazy load Whisper model"""
        if self.whisper_model is None:
            try:
                import whisper  # Lazy import to avoid PyTorch DLL issues at startup
                logger.info("Loading Whisper model for transcription...")
                # Use base model for balance between speed and accuracy
                self.whisper_model = whisper.load_model("base")
            except ImportError as e:
                raise ImportError(
                    "Whisper is not installed. Install it with: pip install openai-whisper"
                ) from e
            except Exception as e:
                raise Exception(
                    f"Failed to load Whisper model. This may be due to missing PyTorch dependencies. "
                    f"Error: {e}. Try reinstalling PyTorch: pip install torch --index-url https://download.pytorch.org/whl/cpu"
                ) from e
        return self.whisper_model
    
    def process_document(self, file_path: Path, mime_type: str = None) -> Dict[str, Any]:
        """Process a single document file"""
        log_activity("Document processing started", {"file_path": str(file_path), "mime_type": mime_type})
        
        try:
            suffix = file_path.suffix.lower()
            
            if suffix == '.docx':
                return self._process_docx(file_path)
            elif suffix == '.doc':
                # Try to read as docx; fallback to raw text extraction
                return self._process_doc_fallback(file_path)
            elif suffix == '.pdf':
                return self._process_pdf(file_path)
            elif suffix == '.pptx':
                return self._process_pptx(file_path)
            elif suffix == '.zip':
                # Detect SCORM vs xAPI vs generic ZIP
                return self._process_zip_smart(file_path)
            elif suffix in ['.mp3', '.wav', '.m4a', '.ogg', '.flac', '.aac']:
                return self._process_audio(file_path)
            elif suffix in ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm']:
                return self._process_video(file_path)
            else:
                raise ValueError(f"Unsupported file type: {suffix}")
        except Exception as error:
            logger.error(f"Failed to process document: {error}", exc_info=True)
            raise
    
    def _process_audio(self, file_path: Path) -> Dict[str, Any]:
        """Process audio file and transcribe to text"""
        try:
            from pydub import AudioSegment  # Lazy import
            
            logger.info(f"Transcribing audio file: {file_path}")
            
            # Load Whisper model
            model = self._get_whisper_model()
            
            # Transcribe audio
            # Fix L1: removed language="en" — Whisper now auto-detects language
            result = model.transcribe(str(file_path))
            transcript = result["text"]
            
            # Get audio duration
            audio = AudioSegment.from_file(str(file_path))
            duration_seconds = len(audio) / 1000.0
            
            return {
                "content": transcript,
                "metadata": {
                    "type": "audio",
                    "file_format": file_path.suffix.lower(),
                    "duration_seconds": duration_seconds,
                    "word_count": len(transcript.split()),
                    "language": result.get("language", "en")
                }
            }
        except Exception as error:
            raise Exception(f"Failed to process audio: {error}")
    
    def _process_video(self, file_path: Path) -> Dict[str, Any]:
        """Process video file: extract audio and transcribe"""
        try:
            from moviepy import VideoFileClip  # Lazy import
            
            logger.info(f"Processing video file: {file_path}")
            
            # Extract audio from video
            temp_audio_path = Path(tempfile.mkdtemp()) / f"extracted_audio_{file_path.stem}.wav"
            temp_audio_path.parent.mkdir(parents=True, exist_ok=True)
            
            video = None
            try:
                # Extract audio using moviepy
                video = VideoFileClip(str(file_path))
                video_duration = video.duration
                video_resolution = f"{video.w}x{video.h}" if hasattr(video, 'w') and hasattr(video, 'h') else None
                
                audio = video.audio
                audio.write_audiofile(str(temp_audio_path), verbose=False, logger=None)
                audio.close()
                video.close()
                video = None
                
                # Transcribe the extracted audio
                model = self._get_whisper_model()
                # Fix L1: auto-detect language (do not force English)
                result = model.transcribe(str(temp_audio_path))
                transcript = result["text"]
                
                return {
                    "content": transcript,
                    "metadata": {
                        "type": "video",
                        "file_format": file_path.suffix.lower(),
                        "duration_seconds": video_duration,
                        "word_count": len(transcript.split()),
                        "language": result.get("language", "en"),
                        "video_resolution": video_resolution
                    }
                }
            finally:
                # Cleanup temp audio file and video
                if video:
                    try:
                        video.close()
                    except:
                        pass
                if temp_audio_path.exists():
                    try:
                        os.remove(temp_audio_path)
                    except:
                        pass
        except Exception as error:
            raise Exception(f"Failed to process video: {error}")
    
    def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """Process DOCX file — preserves reading order including table cells"""
        try:
            from docx.oxml.ns import qn
            doc = Document(file_path)
            content = []

            # Iterate the document body in XML order to preserve paragraph+table order
            for block in doc.element.body:
                tag = block.tag.split('}')[-1] if '}' in block.tag else block.tag
                if tag == 'p':
                    # Paragraph
                    from docx.text.paragraph import Paragraph
                    para = Paragraph(block, doc)
                    if para.text.strip():
                        content.append(para.text.strip())
                elif tag == 'tbl':
                    # Table — iterate rows/cells in order
                    from docx.table import Table
                    table = Table(block, doc)
                    for row in table.rows:
                        row_text = ' | '.join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            content.append(row_text)

            full_content = '\n\n'.join(content)
            return {
                "content": full_content,
                "metadata": {
                    "type": "docx",
                    "word_count": len(full_content.split()),
                    "paragraphs": len(doc.paragraphs)
                }
            }
        except Exception as error:
            raise Exception(f"Failed to process DOCX: {error}")
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Process PDF file with multi-tier extraction for best multi-script support.
        Tries: pdfplumber (best for Hindi/Indic/CJK) → PyMuPDF → PyPDF2 (fallback)
        """
        # --- Tier 1: pdfplumber (best Devanagari / multi-script layout support) ---
        try:
            import pdfplumber
            content = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text(x_tolerance=3, y_tolerance=3) or ""
                    if text.strip():
                        content.append(text.strip())
            if content:
                full_content = '\n\n'.join(content)
                if not _is_text_garbled(full_content):  # quality guard: too-short or null-byte-heavy text means custom font
                    logger.info(f"PDF extracted via pdfplumber: {file_path.name} ({len(content)} pages)")
                    return {
                        "content": full_content,
                        "metadata": {
                            "type": "pdf",
                            "extractor": "pdfplumber",
                            "pages": len(content),
                            "word_count": len(full_content.split())
                        }
                    }
        except ImportError:
            logger.debug("pdfplumber not installed — trying PyMuPDF")
        except Exception as e:
            logger.warning(f"pdfplumber failed for {file_path.name}: {e} — trying PyMuPDF")

        # --- Tier 2: PyMuPDF / fitz (excellent Unicode + font rendering) ---
        try:
            import fitz  # PyMuPDF
            content = []
            with fitz.open(str(file_path)) as doc:
                for page in doc:
                    # 'text' layout mode preserves reading order
                    text = page.get_text("text").strip()
                    if text:
                        content.append(text)
            if content:
                full_content = '\n\n'.join(content)
                if not _is_text_garbled(full_content):
                    logger.info(f"PDF extracted via PyMuPDF: {file_path.name} ({len(content)} pages)")
                    return {
                        "content": full_content,
                        "metadata": {
                            "type": "pdf",
                            "extractor": "pymupdf",
                            "pages": len(content),
                            "word_count": len(full_content.split())
                        }
                    }
        except ImportError:
            logger.debug("PyMuPDF not installed — trying PyPDF2")
        except Exception as e:
            logger.warning(f"PyMuPDF failed for {file_path.name}: {e} — trying PyPDF2")

        # --- Tier 3: PyPDF2 (legacy fallback) ---
        try:
            reader = PdfReader(file_path)
            content = []
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    content.append(text.strip())
            full_content = '\n\n'.join(content)
            if not _is_text_garbled(full_content):
                logger.info(f"PDF extracted via PyPDF2: {file_path.name} ({len(content)} pages)")
                return {
                    "content": full_content,
                    "metadata": {
                        "type": "pdf",
                        "extractor": "pypdf2",
                        "pages": len(reader.pages),
                        "word_count": len(full_content.split())
                    }
                }
            logger.warning(f"PyPDF2 text appears garbled for {file_path.name} — escalating to Vision OCR")
        except Exception as error:
            logger.warning(f"PyPDF2 also failed for {file_path.name}: {error}")

        # --- Tier 4: Gemini Vision OCR (handles custom-encoded / image-based PDFs) ---
        try:
            return self._extract_pdf_via_gemini_vision(file_path)
        except Exception as e:
            raise Exception(f"All extractors including Vision OCR failed for {file_path.name}: {e}")

    def _extract_pdf_via_gemini_vision(self, file_path: Path) -> Dict[str, Any]:
        """Last-resort: render PDF pages as images and use Gemini Vision to read them.
        Works for ANY font encoding including custom Type3 Devanagari, CJK, Arabic fonts.
        Caps at first 5 pages to control cost and latency.
        """
        import fitz  # PyMuPDF required for page rendering
        import base64
        import google.generativeai as genai
        import json
        from config import GEMINI_API_KEY, GEMINI_MODEL

        genai.configure(api_key=GEMINI_API_KEY)
        model_name = GEMINI_MODEL or "gemini-1.5-pro-latest"
        model = genai.GenerativeModel(model_name)

        doc = fitz.open(str(file_path))
        total_pages = len(doc)
        max_pages = min(total_pages, 5)  # cap at 5 pages to avoid excessive tokens
        logger.info(f"PDF Vision OCR: rendering {max_pages}/{total_pages} pages for {file_path.name}")

        page_images = []
        for page_num in range(max_pages):
            page = doc[page_num]
            # Render at 150 DPI (good OCR quality, reasonable size)
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("jpeg", jpg_quality=85)
            page_images.append(img_bytes)
        doc.close()

        # Build a multimodal prompt with all page images
        prompt_parts = [
            "You are an expert OCR assistant. Read ALL text from the provided PDF page image(s) "
            "exactly as it appears, in the original language. "
            "Preserve paragraph structure. Return ONLY the extracted text, nothing else."
        ]
        for i, img_bytes in enumerate(page_images):
            prompt_parts.append({"mime_type": "image/jpeg", "data": img_bytes})

        response = model.generate_content(
            prompt_parts,
            generation_config={"temperature": 0.0, "max_output_tokens": 8192}
        )
        extracted = response.text.strip() if response.text else ""
        words = extracted.split()
        logger.info(
            f"PDF Vision OCR complete: {file_path.name} — {len(words)} words extracted from {max_pages} pages"
        )
        return {
            "content": extracted,
            "metadata": {
                "type": "pdf",
                "extractor": "gemini-vision",
                "pages": max_pages,
                "total_pages": total_pages,
                "word_count": len(words),
                "note": "Vision OCR used — custom/Indic font detected in text layer"
            }
        }

    def _process_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Process PPTX file"""
        try:
            prs = Presentation(file_path)
            content = []
            slides = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slide_content = '\n'.join(slide_text)
                    slides.append(f"Slide {slide_num + 1}:\n{slide_content}")
                    content.append(slide_content)
            
            full_content = '\n\n---\n\n'.join(slides)
            
            return {
                "content": full_content,
                "metadata": {
                    "type": "pptx",
                    "slides": len(prs.slides),
                    "word_count": len(full_content.split())
                }
            }
        except Exception as error:
            raise Exception(f"Failed to process PPTX: {error}")
    
    def _process_zip_smart(self, file_path: Path) -> Dict[str, Any]:
        """Auto-detect ZIP type: SCORM 1.2, xAPI, or generic — then dispatch appropriately"""
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                names = z.namelist()
            
            # SCORM 1.2 detection: has imsmanifest.xml
            if any('imsmanifest.xml' in n for n in names):
                return self._process_scorm_zip(file_path)
            
            # xAPI detection: has tincan.xml or course.json
            if any('tincan.xml' in n for n in names) or any('course.json' in n for n in names):
                return self._process_xapi_zip(file_path)
            
            # Fallback: generic ZIP
            return self._process_zip(file_path)
        except Exception as error:
            raise Exception(f"Failed to detect ZIP type: {error}")

    def _process_scorm_zip(self, file_path: Path) -> Dict[str, Any]:
        """Process SCORM 1.2 ZIP package: parse imsmanifest.xml + extract HTML slide text"""
        import xml.etree.ElementTree as ET
        from html.parser import HTMLParser

        class HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text_parts = []
                self._skip_tags = {'script', 'style'}
                self._current_skip = False
            def handle_starttag(self, tag, attrs):
                if tag.lower() in self._skip_tags:
                    self._current_skip = True
            def handle_endtag(self, tag):
                if tag.lower() in self._skip_tags:
                    self._current_skip = False
            def handle_data(self, data):
                if not self._current_skip and data.strip():
                    self.text_parts.append(data.strip())

        temp_dir = Path(tempfile.mkdtemp())
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                z.extractall(temp_dir)

            # Parse manifest
            manifest_path = next(temp_dir.rglob('imsmanifest.xml'), None)
            course_title = file_path.stem
            sco_files = []

            if manifest_path:
                try:
                    tree = ET.parse(manifest_path)
                    root = tree.getroot()
                    # Strip namespaces for simple access
                    ns = re.match(r'\{.*\}', root.tag)
                    ns_str = ns.group(0) if ns else ''

                    # Course title from organizations
                    for org in root.iter(f'{ns_str}organization'):
                        title_el = org.find(f'{ns_str}title')
                        if title_el is not None and title_el.text:
                            course_title = title_el.text.strip()
                            break

                    # SCO resource hrefs
                    for resource in root.iter(f'{ns_str}resource'):
                        href = resource.get('href')
                        scorm_type = resource.get('adlcp:scormtype') or resource.get('scormtype', '')
                        if href and ('sco' in scorm_type.lower() or href.endswith('.html') or href.endswith('.htm')):
                            sco_files.append(href.split('?')[0])  # Strip query params
                except Exception as e:
                    logger.warning(f"SCORM manifest parse warning: {e}")

            # Extract text from HTML SCO files
            extracted_parts = [f"SCORM Course: {course_title}\n"]
            html_files = list(temp_dir.rglob('*.html')) + list(temp_dir.rglob('*.htm'))

            # Prioritize SCO files found in manifest, else all HTML
            target_files = []
            for sf in sco_files:
                candidate = temp_dir / sf
                if candidate.exists():
                    target_files.append(candidate)
            if not target_files:
                target_files = html_files[:20]  # Cap at 20 to avoid huge files

            for html_file in target_files:
                try:
                    html_content = html_file.read_text(encoding='utf-8', errors='ignore')
                    parser = HTMLTextExtractor()
                    parser.feed(html_content)
                    text = ' '.join(parser.text_parts)
                    if text.strip():
                        extracted_parts.append(f"--- {html_file.name} ---\n{text}")
                except Exception as e:
                    logger.warning(f"Failed to extract text from {html_file.name}: {e}")

            full_content = '\n\n'.join(extracted_parts)
            return {
                "content": full_content,
                "metadata": {
                    "type": "scorm_zip",
                    "course_title": course_title,
                    "sco_count": len(target_files),
                    "word_count": len(full_content.split())
                }
            }
        finally:
            if temp_dir.exists():
                shutil.rmtree(temp_dir)

    def _process_xapi_zip(self, file_path: Path) -> Dict[str, Any]:
        """Process xAPI ZIP: handles Kartavya course.json, Articulate Rise 360, and generic tincan.xml formats"""
        import xml.etree.ElementTree as ET
        import json
        import base64

        temp_dir = Path(tempfile.mkdtemp())
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                z.extractall(temp_dir)

            course_title = file_path.stem
            content_parts = []
            module_titles = []

            # ── Strategy 1: Kartavya own course.json ──────────────────────────
            for cj in temp_dir.rglob('course.json'):
                try:
                    data = json.loads(cj.read_text(encoding='utf-8'))
                    if data.get('course', {}).get('title'):
                        course_title = data['course']['title']
                    content_parts.append(
                        f"[Existing Course JSON]\n{json.dumps(data, indent=2, ensure_ascii=False)[:30000]}"
                    )
                    logger.info(f"Extracted Kartavya course.json: {course_title}")
                    break
                except Exception as e:
                    logger.warning(f"Failed to read course.json: {e}")

            if content_parts:
                # Already handled — skip other strategies
                full_content = '\n\n'.join(content_parts)
                return {
                    "content": full_content,
                    "metadata": {
                        "type": "xapi_zip",
                        "course_title": course_title,
                        "word_count": len(full_content.split())
                    }
                }

            # ── Strategy 2: Articulate Rise 360 (base64 JSON in index.html) ───
            # Rise 360 embeds ALL course data as one huge base64 JSON string inside HTML
            rise_index = temp_dir / 'scormcontent' / 'index.html'
            if rise_index.exists():
                try:
                    html_content = rise_index.read_text(encoding='utf-8', errors='ignore')

                    # Find continuous base64 strings >200 chars (the course data blob)
                    b64_candidates = re.findall(r'[A-Za-z0-9+/]{200,}={0,2}', html_content)

                    rise_data = None
                    for candidate in b64_candidates:
                        try:
                            decoded = base64.b64decode(candidate + '==').decode('utf-8', errors='ignore')
                            if '"lessons"' in decoded and '"course"' in decoded:
                                rise_data = json.loads(decoded)
                                break
                        except Exception:
                            continue

                    if rise_data:
                        logger.info("Detected Articulate Rise 360 xAPI format — extracting content")

                        # Extract course title
                        course_title = (
                            rise_data.get('course', {}).get('title')
                            or rise_data.get('title')
                            or course_title
                        )

                        extracted_text = [f"Course Title: {course_title}\n"]

                        # Extract lessons (= modules in Rise)
                        lessons = rise_data.get('course', {}).get('lessons', []) or rise_data.get('lessons', [])
                        extracted_text.append(f"Total Modules: {len(lessons)}\n")

                        for lesson_idx, lesson in enumerate(lessons, 1):
                            lesson_title = lesson.get('title', f'Module {lesson_idx}')
                            module_titles.append(lesson_title)
                            extracted_text.append(f"\n=== Module {lesson_idx}: {lesson_title} ===")

                            lesson_desc = lesson.get('description', '')
                            if lesson_desc:
                                extracted_text.append(lesson_desc)

                            # Extract text from Rise block items recursively
                            items = lesson.get('items', [])
                            lesson_text = self._extract_rise_block_text(items)
                            if lesson_text:
                                extracted_text.append(lesson_text)

                        full_content = '\n'.join(extracted_text)
                        logger.info(
                            f"Rise 360 extraction complete: {len(lessons)} modules, "
                            f"{len(full_content.split())} words"
                        )
                        return {
                            "content": full_content,
                            "metadata": {
                                "type": "xapi_rise360",
                                "course_title": course_title,
                                "module_count": len(lessons),
                                "module_titles": module_titles,
                                "word_count": len(full_content.split())
                            }
                        }
                    else:
                        logger.warning("Rise 360 index.html found but could not decode course data blob")
                except Exception as e:
                    logger.warning(f"Rise 360 extraction failed: {e}", exc_info=True)

            # ── Strategy 3: tincan.xml — parse activities as module list ──────
            for tc in temp_dir.rglob('tincan.xml'):
                try:
                    tree = ET.parse(tc)
                    root = tree.getroot()
                    # Remove namespace prefix for easier parsing
                    ns_map = {'tc': 'http://projecttincan.com/tincan.xsd'}

                    # Get course title from first "course" type activity
                    for activity in root.iter('{http://projecttincan.com/tincan.xsd}activity'):
                        act_type = activity.get('type', '')
                        name_el = activity.find('{http://projecttincan.com/tincan.xsd}name')
                        if name_el is not None and name_el.text:
                            if 'course' in act_type.lower():
                                course_title = name_el.text.strip()
                            elif 'module' in act_type.lower() or 'lesson' in act_type.lower():
                                module_titles.append(name_el.text.strip())

                    # Build content from tincan activities
                    content_text = [f"Course Title: {course_title}"]
                    if module_titles:
                        content_text.append(f"\nDetected {len(module_titles)} modules:")
                        for idx, mt in enumerate(module_titles, 1):
                            content_text.append(f"  Module {idx}: {mt}")

                    content_parts.append('\n'.join(content_text))
                    logger.info(f"tincan.xml: {course_title} ({len(module_titles)} modules)")
                    break
                except Exception as e:
                    logger.warning(f"Failed to read tincan.xml: {e}")

            # ── Strategy 4: Fallback — extract text from any HTML in the ZIP ──
            if not content_parts:
                from html.parser import HTMLParser

                class _TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.parts = []
                        self._skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag.lower() in ('script', 'style'):
                            self._skip = True
                    def handle_endtag(self, tag):
                        if tag.lower() in ('script', 'style'):
                            self._skip = False
                    def handle_data(self, data):
                        if not self._skip and data.strip():
                            self.parts.append(data.strip())

                html_files = list(temp_dir.rglob('*.html')) + list(temp_dir.rglob('*.htm'))
                for hf in html_files[:5]:
                    try:
                        parser = _TextExtractor()
                        parser.feed(hf.read_text(encoding='utf-8', errors='ignore'))
                        txt = ' '.join(parser.parts)
                        if txt.strip():
                            content_parts.append(f"[{hf.name}]\n{txt[:5000]}")
                    except Exception:
                        pass

            full_content = '\n\n'.join(content_parts) if content_parts else f"xAPI package: {course_title}"
            return {
                "content": full_content,
                "metadata": {
                    "type": "xapi_zip",
                    "course_title": course_title,
                    "module_titles": module_titles,
                    "word_count": len(full_content.split())
                }
            }
        finally:
            if temp_dir.exists():
                shutil.rmtree(temp_dir)

    def _extract_rise_block_text(self, items: list, depth: int = 0) -> str:
        """Recursively extract plain text from Articulate Rise 360 block items"""
        if not items or depth > 5:
            return ''
        parts = []
        for item in items:
            if not isinstance(item, dict):
                continue
            item_type = item.get('type', '')

            # Text/paragraph blocks
            for text_key in ('text', 'content', 'body', 'caption', 'question', 'answer',
                             'title', 'heading', 'label', 'description', 'feedback'):
                val = item.get(text_key)
                if isinstance(val, str) and val.strip():
                    # Strip HTML tags from Rise's rich text fields
                    clean = re.sub(r'<[^>]+>', ' ', val).strip()
                    clean = re.sub(r'\s+', ' ', clean)
                    if clean and len(clean) > 2:
                        parts.append(clean)

            # Choices / options in quiz blocks
            choices = item.get('choices') or item.get('options') or []
            if isinstance(choices, list):
                for choice in choices:
                    if isinstance(choice, dict):
                        for ck in ('text', 'label', 'body'):
                            cv = choice.get(ck)
                            if isinstance(cv, str) and cv.strip():
                                clean = re.sub(r'<[^>]+>', ' ', cv).strip()
                                if clean:
                                    parts.append(f"  • {clean}")

            # Recurse into nested items
            child_items = item.get('items') or item.get('blocks') or item.get('slides') or []
            if isinstance(child_items, list):
                child_text = self._extract_rise_block_text(child_items, depth + 1)
                if child_text:
                    parts.append(child_text)

        return '\n'.join(parts)


    def _process_doc_fallback(self, file_path: Path) -> Dict[str, Any]:
        """Try to read legacy .doc file — attempt docx first, then raw bytes"""
        try:
            # Some .doc files are actually DOCX
            doc = Document(file_path)
            content = [p.text for p in doc.paragraphs if p.text.strip()]
            full_content = '\n\n'.join(content)
            return {
                "content": full_content,
                "metadata": {"type": "doc", "word_count": len(full_content.split())}
            }
        except Exception:
            # Fallback: return filename-based stub
            return {
                "content": f"Legacy .doc file: {file_path.stem}. Content extraction not fully supported for binary .doc format.",
                "metadata": {"type": "doc_legacy", "word_count": 0}
            }

    def _process_zip(self, file_path: Path) -> Dict[str, Any]:
        """Process generic ZIP file containing multiple documents"""
        try:
            temp_dir = Path(tempfile.mkdtemp())
            processed_files = []
            
            try:
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                
                # Process all supported files
                supported_extensions = ['.docx', '.pdf', '.pptx', '.mp3', '.wav', '.m4a', '.mp4', '.avi', '.mov']
                
                for file in temp_dir.rglob('*'):
                    if file.is_file() and file.suffix.lower() in supported_extensions:
                        try:
                            processor = DocumentProcessor()
                            result = processor.process_document(file)
                            processed_files.append({
                                "filename": file.name,
                                "content": result["content"],
                                "metadata": result["metadata"]
                            })
                        except Exception as e:
                            logger.warning(f"Failed to process {file.name}: {e}")
                
                # Combine all content
                combined_content = '\n\n---DOCUMENT SEPARATOR---\n\n'.join(
                    [f"File: {f['filename']}\n\n{f['content']}" for f in processed_files]
                )
                
                return {
                    "content": combined_content,
                    "metadata": {
                        "type": "zip",
                        "files_processed": len(processed_files),
                        "files": [f["filename"] for f in processed_files]
                    },
                    "processed_files": processed_files
                }
            finally:
                # Cleanup temp directory
                if temp_dir.exists():
                    shutil.rmtree(temp_dir)
        except Exception as error:
            raise Exception(f"Failed to process ZIP: {error}")
    
    def analyze_document(self, content: str) -> Dict[str, Any]:
        """Analyze document content for course generation"""
        words = content.split()
        word_count = len(words)
        
        # Simple topic extraction (can be enhanced)
        sentences = [s.strip() for s in content.split('.') if len(s.strip()) > 20]
        topics = sentences[:10]
        
        return {
            "word_count": word_count,
            "estimated_time_minutes": max(10, word_count // 200),  # Rough estimate
            "topics": topics
        }

